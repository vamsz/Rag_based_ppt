import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement

class AdvancedDesignEngine:
    def __init__(self):
        # 1. Define a "Cyberpunk/Glass" Color Palette
        self.colors = {
            "bg": RGBColor(10, 25, 47),          # Dark Navy Background
            "glass": RGBColor(23, 42, 69),       # Lighter Navy for "Glass" cards
            "accent": RGBColor(100, 255, 218),   # Neon Cyan/Green
            "text_main": RGBColor(230, 241, 255),# Off-white
            "text_dim": RGBColor(136, 146, 176)  # Muted Grey
        }

    def create_presentation(self, slides_data):
        """Main entry point to generate the deck."""
        prs = Presentation()
        # Force 16:9 Widescreen
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        for slide_info in slides_data:
            slide_type = slide_info.get('type', 'content').lower()
            if slide_type == 'title':
                self.render_title_slide(prs, slide_info)
            else:  # content or any other type
                self.render_content_slide(prs, slide_info)
        
        # Return BytesIO for FastAPI streaming
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output

    def render_title_slide(self, prs, data):
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # 6 = Blank Layout
        self._set_background(slide)

        # ADVANCED: Add a "Glowing Orb" behind the text using transparency
        self._add_glass_shape(slide, MSO_SHAPE.OVAL, 8, -2, 6, 6, transparency=80)

        # Draw Title Text
        title = data.get('title', 'Untitled')
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(10), Inches(2))
        self._write_text(title_box, title, 60, True, self.colors['text_main'])

        # Draw Subtitle Text (use content if no subtitle provided)
        subtitle = data.get('subtitle', data.get('content', ''))
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(10), Inches(1))
        self._write_text(sub_box, subtitle, 24, False, self.colors['accent'])

    def render_content_slide(self, prs, data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_background(slide)

        # 1. Header Area
        title = data.get('title', 'Slide')
        self._write_text(
            slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1)),
            title, 36, True, self.colors['accent']
        )

        # 2. The "Glass Card" Container
        # We draw a rounded rectangle behind the text to make it pop
        card = self._add_glass_shape(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE, 
            0.5, 1.8, 12.3, 5.0, transparency=10
        )

        # 3. Content Bullets
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4.5))
        tf = content_box.text_frame
        tf.word_wrap = True

        # Handle both 'points' (list) and 'content' (string or list) formats
        points = data.get('points', [])
        if not points:
            content = data.get('content', '')
            if content:
                # Check if content is already a list
                if isinstance(content, list):
                    points = content
                else:
                    # If it's a string, split by newlines
                    points = [line.strip() for line in content.split('\n') if line.strip()]
        
        for point in points:
            p = tf.add_paragraph()
            p.text = f"• {point}"
            p.space_after = Pt(14)
            # Apply formatting to the paragraph run
            if p.runs:
                run = p.runs[0]
                run.font.size = Pt(20)
                run.font.color.rgb = self.colors['text_main']
                run.font.name = "Arial"

    # --- THE "ADVANCED" HELPERS ---

    def _set_background(self, slide):
        """Paint the background solid dark navy."""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['bg']

    def _add_glass_shape(self, slide, shape_type, x, y, w, h, transparency=0):
        """
        Draws a shape and safely attempts to set transparency.
        """
        shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.colors['glass']
        shape.line.fill.background() # Remove outline

        if transparency > 0:
            try:
                # THIS IS THE ADVANCED PART: XML MANIPULATION
                alpha_value = int((100 - transparency) * 1000)
                
                # Safe XML access pattern for python-pptx
                fill = shape.fill
                # Get the solidFill element (either locally or from the theme)
                if hasattr(fill, '_xPr'):
                    fill_element = fill._xPr.solidFill
                else:
                    # Fallback or alternative access
                    return shape

                if fill_element is not None:
                    # srgbClr might be the color type used
                    color_element = fill_element.srgbClr
                    if color_element is not None:
                        # Inject the alpha tag
                        alpha = OxmlElement('a:alpha')
                        alpha.set('val', str(alpha_value))
                        color_element.append(alpha)
            except Exception as e:
                print(f"Warning: Could not apply transparency effect: {e}")
                # Do not crash, just continue with solid color
        
        return shape

    def _write_text(self, textbox, text, size, bold, color):
        try:
            tf = textbox.text_frame
            # Ensure there's a paragraph
            if not tf.paragraphs:
                p = tf.add_paragraph()
            else:
                p = tf.paragraphs[0]
            
            p.text = str(text) if text else ""
            
            # Ensure there's a run to style
            if not p.runs:
                p.add_run()
                
            run = p.runs[0]
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = "Arial"
        except Exception as e:
            print(f"Warning: Could not style text '{text}': {e}")

# --- USAGE EXAMPLE (What you asked for) ---

if __name__ == "__main__":
    # 1. Initialize the Engine
    engine = AdvancedDesignEngine()

    # 2. Define your content (This would come from the LLM/RAG usually)
    my_presentation_data = [
        {
            "type": "title",
            "title": "FUTURE OF AI",
            "subtitle": "Generated by Python Design Engine"
        },
        {
            "type": "content",
            "title": "Why this code is Advanced",
            "points": [
                "It ignores standard templates and draws everything from scratch.",
                "It uses 'XML Hacking' to enable Transparency (Glassmorphism).",
                "It separates Data (JSON) from Design (Python Class).",
                "The font sizes and colors are enforced strictly."
            ]
        }
    ]

    # 3. Generate and Save
    print("Generating presentation...")
    ppt_object = engine.create_presentation(my_presentation_data)
    ppt_object.save("Design_Showcase.pptx")
    print("Done! Open 'Design_Showcase.pptx' to see the result.")
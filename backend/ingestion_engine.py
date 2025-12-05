import logging
import os
import tempfile
from typing import Callable

import pandas as pd
from docx import Document
from pptx import Presentation
from pypdf import PdfReader

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class UniversalLoader:
    """
    Universal file loader for extracting text from various file formats.
    """
    def __init__(self) -> None:
        self._loaders: dict[str, Callable[[str], str]] = {
            ".pdf": self._load_pdf,
            ".docx": self._load_docx,
            ".pptx": self._load_pptx,
            ".csv": self._load_csv,
            ".xlsx": self._load_xlsx,
            ".txt": self._load_txt,
        }

    def process_upload(self, uploaded_file) -> str:
        """
        Main entry point. Saves the file temporarily and routes to the correct loader.
        
        Args:
            uploaded_file: Streamlit UploadedFile object.
            
        Returns:
            str: Extracted text content.
        """
        if uploaded_file is None:
            return ""

        file_ext = os.path.splitext(uploaded_file.name)[1].lower().strip()
        
        # Create a temporary file to store the uploaded content
        with tempfile.NamedTemporaryFile(delete=False, suffix=file_ext) as tmp_file:
            tmp_file.write(uploaded_file.getbuffer())
            tmp_path = tmp_file.name

        text = ""
        try:
            logger.info(f"Processing file: {uploaded_file.name} with extension {file_ext}")
            loader = self._loaders.get(file_ext)
            if not loader:
                logger.warning(f"Unsupported file format: {file_ext}")
                return f"Error: Unsupported file format {file_ext}"

            text = loader(tmp_path)
                
        except Exception as e:
            logger.error(f"Error processing file {uploaded_file.name}: {str(e)}")
            return f"Error processing file: {str(e)}"
        finally:
            # Clean up temp file
            if os.path.exists(tmp_path):
                os.remove(tmp_path)

        return text

    def _load_pdf(self, file_path: str) -> str:
        try:
            reader = PdfReader(file_path)
            text = []
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted:
                    text.append(extracted)
            return "\n".join(text)
        except Exception as e:
            logger.error(f"PDF extraction error: {e}")
            raise

    def _load_docx(self, file_path: str) -> str:
        try:
            doc = Document(file_path)
            text = []
            for para in doc.paragraphs:
                if para.text.strip():
                    text.append(para.text)
            
            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text for cell in row.cells]
                    text.append(" | ".join(row_text))
            
            return "\n".join(text)
        except Exception as e:
            logger.error(f"DOCX extraction error: {e}")
            raise

    def _load_pptx(self, file_path: str) -> str:
        try:
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        shape_text = shape.text_frame.text.strip()
                        if shape_text:
                            text.append(shape_text)
            return "\n".join(text)
        except Exception as e:
            logger.error(f"PPTX extraction error: {e}")
            raise

    def _load_csv(self, file_path: str) -> str:
        try:
            df = pd.read_csv(file_path, on_bad_lines="skip", engine="python")
            if df.empty:
                return ""

            rows = []
            for _, row in df.iterrows():
                row_str = ", ".join(
                    f"{col}: {val}" for col, val in row.items() if pd.notna(val)
                )
                if row_str:
                    rows.append(row_str)
            return "\n".join(rows)
        except Exception as e:
            logger.error(f"CSV extraction error: {e}")
            # Fallback to raw read to salvage text from severely malformed files.
            try:
                with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                    return f.read()
            except Exception:
                raise

    def _load_xlsx(self, file_path: str) -> str:
        try:
            # Read all sheets
            dfs = pd.read_excel(file_path, sheet_name=None, engine="openpyxl")
            text = []
            for sheet_name, df in dfs.items():
                text.append(f"Sheet: {sheet_name}")
                for index, row in df.iterrows():
                    row_str = ", ".join([f"{col}: {val}" for col, val in row.items() if pd.notna(val)])
                    text.append(row_str)
            return "\n".join(text)
        except Exception as e:
            logger.error(f"Excel extraction error: {e}")
            raise

    def _load_txt(self, file_path: str) -> str:
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                return f.read()
        except Exception as e:
            logger.error(f"TXT extraction error: {e}")
            raise

